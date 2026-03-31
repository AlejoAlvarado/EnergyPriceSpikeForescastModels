from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
PROJECT_DIR = BASE_DIR.parent.parent
ASSET_DIR = BASE_DIR / "assets"

MODELING_DATA = PROJECT_DIR / "JorgeFolder" / "outputs" / "data" / "modeling_dataset.csv"
SPLIT_SUMMARY = PROJECT_DIR / "JorgeFolder" / "outputs" / "data" / "split_summary.csv"
FIGURE_DIR = PROJECT_DIR / "JorgeFolder" / "outputs" / "figures"

TIME_SERIES_FIG = FIGURE_DIR / "time_series_overview.png"
PRICE_DIST_FIG = FIGURE_DIR / "price_distribution.png"
PRICE_HEATMAP_FIG = FIGURE_DIR / "price_heatmap_hour_month.png"
SPIKE_COMPARE_FIG = FIGURE_DIR / "spike_vs_non_spike.png"

LOAD_HEATMAP_FIG = ASSET_DIR / "load_heatmap_hour_month.png"
MODEL_SCORE_FIG = ASSET_DIR / "model_scores_summary.png"
OUTPUT_PPTX = BASE_DIR / "alberta_price_spikes_summary_deck.pptx"


# Professional neutral palette with a restrained red accent.
RED = RGBColor(181, 32, 46)
RED_SOFT = RGBColor(247, 233, 236)
INK = RGBColor(36, 36, 36)
SLATE = RGBColor(90, 90, 90)
MIST = RGBColor(243, 244, 246)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(59, 89, 152)
GOLD = RGBColor(191, 144, 0)


def load_data():
    df = pd.read_csv(MODELING_DATA, parse_dates=["datetime"])
    split = pd.read_csv(SPLIT_SUMMARY)
    return df, split


def create_load_heatmap(df: pd.DataFrame):
    tmp = df.copy()
    tmp["month"] = tmp["datetime"].dt.month
    pivot = tmp.pivot_table(values="ACTUAL_AIL", index="month", columns="hour_of_day", aggfunc="mean")

    sns.set_theme(style="white")
    fig, ax = plt.subplots(figsize=(8.2, 4.2), dpi=180)
    sns.heatmap(
        pivot,
        cmap="Blues",
        cbar_kws={"label": "Average AIL (MW)"},
        ax=ax,
        linewidths=0.2,
        linecolor="white",
    )
    ax.set_title("Average Demand by Month and Hour", fontsize=13, weight="bold")
    ax.set_xlabel("Hour of day")
    ax.set_ylabel("Month")
    fig.tight_layout()
    fig.savefig(LOAD_HEATMAP_FIG, bbox_inches="tight")
    plt.close(fig)


def create_model_chart():
    scores = pd.DataFrame(
        {
            "Model": ["MLP", "LSTM", "CNN"],
            "F1": [0.363, 0.394, 0.416],
            "Color": ["#b5202e", "#5f6b7a", "#2f6f5f"],
        }
    ).sort_values("F1")

    sns.set_theme(style="whitegrid")
    fig, ax = plt.subplots(figsize=(7.4, 4.2), dpi=180)
    bars = ax.barh(scores["Model"], scores["F1"], color=scores["Color"])
    ax.set_xlim(0, 0.5)
    ax.set_xlabel("Test F1-score")
    ax.set_title("Model Comparison on the Test Set", fontsize=13, weight="bold")
    ax.grid(axis="x", alpha=0.25)
    ax.grid(axis="y", visible=False)
    for bar, value in zip(bars, scores["F1"]):
        ax.text(value + 0.012, bar.get_y() + bar.get_height() / 2, f"{value:.3f}", va="center", fontsize=10)
    fig.tight_layout()
    fig.savefig(MODEL_SCORE_FIG, bbox_inches="tight")
    plt.close(fig)


def add_slide_number(slide, number: int):
    box = slide.shapes.add_textbox(Inches(12.4), Inches(0.23), Inches(0.4), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(number)
    run.font.name = "Arial"
    run.font.size = Pt(10)
    run.font.color.rgb = SLATE


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.28), Inches(10.8), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = INK

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(1.6), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.02), Inches(11.2), Inches(0.32))
        p = sub.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Arial"
        run.font.size = Pt(10.5)
        run.font.color.rgb = SLATE


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(8)
    run.font.color.rgb = SLATE


def add_bullets(slide, bullets, x, y, w, h, font_size=18):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = INK
        p.bullet = True
        p.space_after = Pt(8)
    return box


def add_callout(slide, title, body, x, y, w, h, fill=RED_SOFT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill

    title_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12), w - Inches(0.3), Inches(0.22))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = INK

    body_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.4), w - Inches(0.32), h - Inches(0.5))
    body_box.text_frame.word_wrap = True
    p = body_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.name = "Arial"
    run.font.size = Pt(11.5)
    run.font.color.rgb = INK


def add_split_timeline(slide, split: pd.DataFrame, x, y, w, h):
    total_rows = split["rows"].sum()
    curr_x = x
    colors = [RGBColor(225, 239, 254), RGBColor(255, 243, 205), RGBColor(231, 244, 234)]
    for idx, row in split.iterrows():
        width = w * (row["rows"] / total_rows)
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, curr_x, y, width, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = colors[idx]
        rect.line.color.rgb = WHITE

        label = slide.shapes.add_textbox(curr_x + Inches(0.08), y + Inches(0.05), width - Inches(0.16), h - Inches(0.1))
        tf = label.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{row['split'].title()}\n{int(row['rows']):,} rows\n{row['spike_rate_lead_2'] * 100:.1f}% spikes"
        run.font.name = "Arial"
        run.font.size = Pt(11)
        run.font.color.rgb = INK
        curr_x += width


def build_deck():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    df, split = load_data()
    create_load_heatmap(df)
    create_model_chart()

    overall_spike_rate = df["spike_lead_2"].mean() * 100
    start = df["datetime"].min()
    end = df["datetime"].max()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(0.12), Inches(4.9))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(1.25), Inches(1.05), Inches(10.8), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Predicting Alberta Electricity Price Spikes"
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = INK
    p2 = tf.add_paragraph()
    run = p2.add_run()
    run.text = "A summary presentation of the EDA, modeling setup, and final results"
    run.font.name = "Arial"
    run.font.size = Pt(18)
    run.font.color.rgb = SLATE

    callout = slide.shapes.add_textbox(Inches(1.25), Inches(3.0), Inches(8.8), Inches(1.4))
    p = callout.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Models summarized here: CNN = 0.416, LSTM = 0.394, MLP = 0.363"
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.color.rgb = RED
    run.font.bold = True

    p = callout.text_frame.add_paragraph()
    run = p.add_run()
    run.text = "Target: predict whether the Alberta pool price exceeds CAD 200/MWh at t+2."
    run.font.name = "Arial"
    run.font.size = Pt(16)
    run.font.color.rgb = INK

    add_callout(
        slide,
        "Project scope",
        "Use public AESO hourly market and generation data to identify short-run spike risk and compare simple and sequential neural models.",
        Inches(9.55),
        Inches(3.05),
        Inches(2.85),
        Inches(1.55),
        fill=MIST,
    )
    add_footer(slide, "Data 607 presentation summary")
    add_slide_number(slide, 1)

    # Slide 2: Motivation and data
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Why This Problem Matters")
    add_bullets(
        slide,
        [
            "Alberta is an energy-only market, so tight supply-demand conditions can translate into sharp hourly price spikes.",
            "Those spikes matter operationally and financially for retailers, generators, large consumers, and risk managers.",
            "The goal is to predict spike risk two hours ahead using only information available at time t.",
        ],
        Inches(0.8),
        Inches(1.55),
        Inches(7.3),
        Inches(3.1),
        font_size=18,
    )
    add_callout(
        slide,
        "Modeling sample",
        f"{len(df):,} modeled hourly rows\n{start:%Y-%m-%d} to {end:%Y-%m-%d %H:%M}\nOverall spike rate: {overall_spike_rate:.1f}%",
        Inches(8.55),
        Inches(1.75),
        Inches(3.85),
        Inches(1.55),
        fill=RED_SOFT,
    )
    add_callout(
        slide,
        "Main inputs",
        "Pool price, load, imports/exports, generation by fuel type, renewable shares, reserve-style tightness indicators, and lagged features.",
        Inches(8.55),
        Inches(3.7),
        Inches(3.85),
        Inches(1.55),
        fill=MIST,
    )
    add_callout(
        slide,
        "Two AESO sources",
        "Hourly Metered Volumes / Pool Price / AIL\n+\nHistorical Generation Data (CSD)",
        Inches(8.55),
        Inches(5.55),
        Inches(3.85),
        Inches(1.0),
        fill=RED_SOFT,
    )
    add_footer(slide, "Source: AESO hourly market and generation datasets")
    add_slide_number(slide, 2)

    # Slide 3: Time series
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Time Series Overview", "Price, load, and generation move together through changing system conditions.")
    slide.shapes.add_picture(str(TIME_SERIES_FIG), Inches(0.75), Inches(1.45), width=Inches(11.8))
    add_callout(
        slide,
        "Reading the chart",
        "The series shows that Alberta prices are calm most of the time, but they jump sharply in stressed hours. Load and generation conditions provide the context for those jumps.",
        Inches(8.85),
        Inches(5.55),
        Inches(3.0),
        Inches(1.05),
        fill=MIST,
    )
    add_footer(slide, "Project EDA figure: time series overview")
    add_slide_number(slide, 3)

    # Slide 4: EDA distribution + spike conditions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "EDA: Heavy Tails And Pre-Spike Conditions")
    slide.shapes.add_picture(str(PRICE_DIST_FIG), Inches(0.7), Inches(1.45), width=Inches(5.75))
    slide.shapes.add_picture(str(SPIKE_COMPARE_FIG), Inches(6.65), Inches(1.45), width=Inches(5.95))
    add_bullets(
        slide,
        [
            "The price distribution is strongly right-skewed, which is why spike classification is more useful than plain average-price forecasting.",
            "Rows followed by spikes already show higher stress: higher current prices and net load, but lower wind output and reserve margin.",
        ],
        Inches(0.95),
        Inches(5.95),
        Inches(11.2),
        Inches(0.9),
        font_size=16,
    )
    add_footer(slide, "Project EDA figures: price distribution and spike vs. non-spike comparison")
    add_slide_number(slide, 4)

    # Slide 5: Hour and month patterns
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "EDA: Patterns By Hour And Month")
    slide.shapes.add_picture(str(PRICE_HEATMAP_FIG), Inches(0.7), Inches(1.45), width=Inches(6.0))
    slide.shapes.add_picture(str(LOAD_HEATMAP_FIG), Inches(6.85), Inches(1.45), width=Inches(5.8))
    add_bullets(
        slide,
        [
            "The price heatmap shows that high-price periods cluster in specific hours and months rather than appearing uniformly at random.",
            "The demand heatmap shows predictable hourly and seasonal usage patterns, which helps explain why some periods are consistently more vulnerable to spikes.",
        ],
        Inches(0.95),
        Inches(5.95),
        Inches(11.2),
        Inches(0.9),
        font_size=16,
    )
    add_footer(slide, "Left: average price by month/hour. Right: average AIL by month/hour.")
    add_slide_number(slide, 5)

    # Slide 6: Modeling setup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Modeling Setup")
    add_callout(
        slide,
        "Target",
        "Binary spike indicator at t+2\n1 if pool price > CAD 200/MWh, else 0",
        Inches(0.8),
        Inches(1.6),
        Inches(2.6),
        Inches(1.25),
        fill=RED_SOFT,
    )
    add_callout(
        slide,
        "Metric",
        "F1-score is the main metric because spikes are rare and precision-recall balance matters more than accuracy.",
        Inches(3.65),
        Inches(1.6),
        Inches(3.25),
        Inches(1.25),
        fill=MIST,
    )
    add_callout(
        slide,
        "Leakage control",
        "Strict time-based splitting, training-only scaling, and untouched test data until final evaluation.",
        Inches(7.15),
        Inches(1.6),
        Inches(5.05),
        Inches(1.25),
        fill=RED_SOFT,
    )
    add_split_timeline(slide, split, Inches(0.9), Inches(3.45), Inches(11.4), Inches(1.0))
    add_bullets(
        slide,
        [
            "Training and threshold selection were done before the final test period.",
            "This makes the reported model comparison more credible as an out-of-sample summary.",
        ],
        Inches(0.95),
        Inches(5.15),
        Inches(11.2),
        Inches(0.9),
        font_size=17,
    )
    add_footer(slide, "Split summary taken from the project outputs")
    add_slide_number(slide, 6)

    # Slide 7: Models
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Three Models, Same Prediction Task")
    add_callout(
        slide,
        "MLP",
        "Uses a flat feature vector with lagged price and spike indicators.\n\nRole in the project: simple neural baseline.",
        Inches(0.8),
        Inches(1.7),
        Inches(3.65),
        Inches(2.2),
        fill=MIST,
    )
    add_callout(
        slide,
        "LSTM",
        "Uses sequential input windows so the model can learn temporal dependence directly.\n\nRole in the project: recurrent benchmark.",
        Inches(4.85),
        Inches(1.7),
        Inches(3.65),
        Inches(2.2),
        fill=RED_SOFT,
    )
    add_callout(
        slide,
        "CNN",
        "Uses temporal filters to detect local ramp and regime patterns in recent observations.\n\nRole in the project: strongest final model.",
        Inches(8.9),
        Inches(1.7),
        Inches(3.65),
        Inches(2.2),
        fill=MIST,
    )
    add_bullets(
        slide,
        [
            "The design intentionally moves from a simple baseline to richer temporal architectures.",
            "That makes it easier to judge whether extra temporal structure actually improves performance.",
        ],
        Inches(0.95),
        Inches(4.6),
        Inches(11.2),
        Inches(1.1),
        font_size=18,
    )
    add_callout(
        slide,
        "What changed across models",
        "The target, split, and evaluation logic stayed fixed. Only the architecture changed.",
        Inches(0.95),
        Inches(6.0),
        Inches(11.0),
        Inches(0.65),
        fill=RED_SOFT,
    )
    add_footer(slide, "Model summary slide for presentation use")
    add_slide_number(slide, 7)

    # Slide 8: Results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Results Summary")
    slide.shapes.add_picture(str(MODEL_SCORE_FIG), Inches(0.75), Inches(1.45), width=Inches(6.2))
    add_callout(
        slide,
        "Ranking",
        "CNN = 0.416\nLSTM = 0.394\nMLP = 0.363",
        Inches(7.55),
        Inches(1.7),
        Inches(2.1),
        Inches(1.35),
        fill=RED_SOFT,
    )
    add_callout(
        slide,
        "Main finding",
        "The CNN delivered the best overall balance of pattern recognition and generalization on the held-out test set.",
        Inches(9.9),
        Inches(1.7),
        Inches(2.4),
        Inches(1.35),
        fill=MIST,
    )
    add_bullets(
        slide,
        [
            "The CNN had the highest test F1, suggesting that local temporal patterns matter for short-horizon spike prediction.",
            "The LSTM remained competitive but did not surpass the CNN.",
            "The MLP served its purpose as a simpler benchmark, but it captured less of the temporal structure in the data.",
        ],
        Inches(7.45),
        Inches(3.55),
        Inches(4.9),
        Inches(2.1),
        font_size=17,
    )
    add_footer(slide, "Scores used in this summary: CNN 0.416, LSTM 0.394, MLP 0.363")
    add_slide_number(slide, 8)

    # Slide 9: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Key Takeaways")
    add_callout(
        slide,
        "Takeaway 1",
        "Electricity prices in Alberta are heavy-tailed, seasonal, and clearly linked to system tightness.",
        Inches(0.85),
        Inches(1.65),
        Inches(3.75),
        Inches(1.45),
        fill=MIST,
    )
    add_callout(
        slide,
        "Takeaway 2",
        "EDA supports the modeling logic: hours before spikes already look different in terms of load, wind, and reserve conditions.",
        Inches(4.8),
        Inches(1.65),
        Inches(3.75),
        Inches(1.45),
        fill=RED_SOFT,
    )
    add_callout(
        slide,
        "Takeaway 3",
        "Among the three neural models summarized here, the CNN produced the strongest final result.",
        Inches(8.75),
        Inches(1.65),
        Inches(3.75),
        Inches(1.45),
        fill=MIST,
    )
    add_bullets(
        slide,
        [
            "This project shows that public AESO data contains useful short-run predictive signal for spike risk.",
            "The summary result is straightforward: CNN first, LSTM second, MLP third.",
            "Next improvements would likely come from weather, outages, and calibration rather than simply adding more text-book complexity.",
        ],
        Inches(0.95),
        Inches(4.0),
        Inches(11.2),
        Inches(1.8),
        font_size=18,
    )
    add_footer(slide, "Professional summary deck generated from project outputs")
    add_slide_number(slide, 9)

    prs.save(OUTPUT_PPTX)
    print(f"Saved deck to {OUTPUT_PPTX}")


if __name__ == "__main__":
    build_deck()

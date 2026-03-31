from pathlib import Path
import json

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
PROJECT_DIR = BASE_DIR.parent.parent
DATA_PATH = PROJECT_DIR / "Data" / "CSVs" / "aeso_merged_2020_2025.csv"
BASELINE_PATH = PROJECT_DIR / "JorgeFolder" / "models" / "baseline" / "baseline_run" / "baseline_comparison.csv"
LSTM_PATH = PROJECT_DIR / "JorgeFolder" / "models" / "lstm" / "random_search_run" / "metrics.json"
MLP_PATH = PROJECT_DIR / "JorgeFolder" / "models" / "mlp" / "random_search_run" / "metrics.json"
ASSET_DIR = BASE_DIR / "assets"
OUTPUT_PPTX = BASE_DIR / "alberta_price_spikes_presentation.pptx"


# Clean academic palette with a Calgary-red accent.
RED = RGBColor(214, 0, 28)
RED_LIGHT = RGBColor(245, 228, 232)
CHARCOAL = RGBColor(44, 44, 44)
MID_GRAY = RGBColor(102, 102, 102)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(31, 119, 80)
BLUE = RGBColor(44, 103, 176)


def load_project_data():
    df = pd.read_csv(DATA_PATH, parse_dates=["datetime"])

    baseline = pd.read_csv(BASELINE_PATH)
    lstm = json.loads(LSTM_PATH.read_text(encoding="utf-8"))
    mlp = json.loads(MLP_PATH.read_text(encoding="utf-8"))

    cnn = {
        "model": "CNN",
        "f1": 0.4157,
        "precision": 0.3397,
        "recall": 0.5354,
        "note": "From executed cnn_v2_colab_postrun.ipynb with validation threshold tuning.",
    }

    return df, baseline, lstm, mlp, cnn


def create_price_histogram(df: pd.DataFrame, out_path: Path):
    sns.set_theme(style="whitegrid")
    fig, ax = plt.subplots(figsize=(8, 4.5), dpi=180)
    ax.hist(df["ACTUAL_POOL_PRICE"], bins=80, color="#2c67b0", alpha=0.85, edgecolor="white")
    ax.axvline(200, color="#d6001c", linestyle="--", linewidth=2, label="Spike threshold = 200")
    ax.set_xlim(0, 1000)
    ax.set_title("Pool Price Distribution, 2020-2025", fontsize=13, weight="bold")
    ax.set_xlabel("CAD/MWh")
    ax.set_ylabel("Hourly count")
    ax.legend(frameon=False, loc="upper right")
    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)


def create_price_heatmap(df: pd.DataFrame, out_path: Path):
    sns.set_theme(style="white")
    pivot = df.pivot_table(values="ACTUAL_POOL_PRICE", index="month", columns="hour_of_day", aggfunc="mean")
    fig, ax = plt.subplots(figsize=(8, 4.8), dpi=180)
    sns.heatmap(
        pivot,
        cmap="YlOrRd",
        cbar_kws={"label": "Average CAD/MWh"},
        ax=ax,
        linewidths=0.2,
        linecolor="white",
    )
    ax.set_title("Average Pool Price by Month and Hour", fontsize=13, weight="bold")
    ax.set_xlabel("Hour of day")
    ax.set_ylabel("Month")
    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)


def create_model_comparison_chart(rows, out_path: Path):
    sns.set_theme(style="whitegrid")
    df = pd.DataFrame(rows).sort_values("f1")
    colors = ["#999999", "#666666", "#d48a96", "#9b1b30", "#d6001c"]
    fig, ax = plt.subplots(figsize=(7.8, 4.4), dpi=180)
    ax.barh(df["model"], df["f1"], color=colors[: len(df)])
    ax.set_xlim(0, 0.5)
    ax.set_xlabel("Test F1-score")
    ax.set_title("Held-out Test F1 by Model", fontsize=13, weight="bold")
    for idx, value in enumerate(df["f1"]):
        ax.text(value + 0.01, idx, f"{value:.3f}", va="center", fontsize=10)
    fig.tight_layout()
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(11.6), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = CHARCOAL

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.05), Inches(1.7), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.12), Inches(11.4), Inches(0.45))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Arial"
        run.font.size = Pt(10.5)
        run.font.color.rgb = MID_GRAY


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(7.0), Inches(12.0), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(8)
    run.font.color.rgb = MID_GRAY


def add_bullets(slide, bullets, x, y, w, h, font_size=18, color=CHARCOAL):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.bullet = True
    return box


def add_callout(slide, title, body, x, y, w, h, fill_color=RED_LIGHT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color

    title_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.12), w - Inches(0.36), Inches(0.25))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = CHARCOAL

    body_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.42), w - Inches(0.36), h - Inches(0.5))
    p = body_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.name = "Arial"
    run.font.size = Pt(11)
    run.font.color.rgb = CHARCOAL
    body_box.text_frame.word_wrap = True


def format_value(label, value):
    if label in {"Current pool price"}:
        return f"{value:,.2f}"
    if label in {"Current demand", "Net load", "Wind output"}:
        return f"{value:,.0f}"
    if label in {"Renewables share"}:
        return f"{100 * value:.1f}%"
    if label in {"Reserve margin"}:
        return f"{value:.3f}"
    return f"{value:,.2f}"


def build_deck():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    df, baseline, lstm, mlp, cnn = load_project_data()

    hist_path = ASSET_DIR / "price_histogram.png"
    heatmap_path = ASSET_DIR / "price_heatmap.png"
    model_chart_path = ASSET_DIR / "model_comparison.png"

    create_price_histogram(df, hist_path)
    create_price_heatmap(df, heatmap_path)

    baseline_rows = baseline.set_index("baseline")
    comparison_rows = [
        {"model": "Naive", "f1": float(baseline_rows.loc["Naive", "test_f1"])},
        {"model": "SARIMAX", "f1": float(baseline_rows.loc["SARIMAX + Fourier", "test_f1"])},
        {"model": "MLP", "f1": float(mlp["test"]["f1"])},
        {"model": "LSTM", "f1": float(lstm["test"]["f1"])},
        {"model": "CNN", "f1": float(cnn["f1"])},
    ]
    create_model_comparison_chart(comparison_rows, model_chart_path)

    spike_rate = df["spike_lead_2"].mean()
    price_stats = df["ACTUAL_POOL_PRICE"].describe(percentiles=[0.5, 0.95, 0.99])
    avg_by_hour = df.groupby("hour_of_day")["ACTUAL_POOL_PRICE"].mean().sort_values(ascending=False)
    avg_by_month = df.groupby("month")["ACTUAL_POOL_PRICE"].mean().sort_values(ascending=False)
    pre_spike = df.groupby("spike_lead_2")[
        ["ACTUAL_POOL_PRICE", "ACTUAL_AIL", "net_load", "wind_total", "renewables_share", "reserve_margin"]
    ].mean()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(1.0), prs.slide_height)
    side.fill.solid()
    side.fill.fore_color.rgb = RED
    side.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.25), Inches(1.0), Inches(10.8), Inches(1.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Predicting Electricity Price Spikes in Alberta"
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = CHARCOAL
    p2 = tf.add_paragraph()
    p2.text = "(2020-2025)"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RED

    subtitle_box = slide.shapes.add_textbox(Inches(1.28), Inches(2.7), Inches(9.8), Inches(1.0))
    p = subtitle_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Data 607 machine learning project\nComparing MLP, LSTM, CNN, and strong time-series baselines"
    run.font.name = "Arial"
    run.font.size = Pt(18)
    run.font.color.rgb = MID_GRAY

    add_callout(
        slide,
        "Core question",
        "Can public AESO market and generation data predict whether the Alberta pool price will exceed CAD 200/MWh two hours ahead?",
        Inches(1.25),
        Inches(4.4),
        Inches(10.4),
        Inches(1.35),
        fill_color=RED_LIGHT,
    )
    add_footer(slide, "Source: AESO data merged and modeled by the project team")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Motivation And Research Question")
    add_bullets(
        slide,
        [
            "Alberta’s deregulated electricity market experiences large hourly pool-price spikes that create financial and operational risk.",
            "Short-horizon spike prediction matters for market participants because dispatch conditions can tighten very quickly.",
            "This project treats spike forecasting as a classification problem instead of only describing past volatility.",
        ],
        Inches(0.75),
        Inches(1.55),
        Inches(7.0),
        Inches(4.6),
    )
    add_callout(
        slide,
        "Research question",
        "Using publicly available AESO data, can we predict whether a spike occurs at t+2 and determine whether sequence models outperform simpler benchmarks?",
        Inches(8.0),
        Inches(1.75),
        Inches(4.6),
        Inches(1.5),
    )
    add_callout(
        slide,
        "Definition of a spike",
        "spike_lead_2 = 1 when the pool price at time t+2 exceeds CAD 200/MWh.",
        Inches(8.0),
        Inches(3.55),
        Inches(4.6),
        Inches(1.15),
        fill_color=LIGHT_GRAY,
    )
    add_callout(
        slide,
        "Why the horizon matters",
        "Using t+2 prevents trivial contemporaneous prediction and makes the exercise closer to a short-run operational warning system.",
        Inches(8.0),
        Inches(4.95),
        Inches(4.6),
        Inches(1.45),
        fill_color=LIGHT_GRAY,
    )
    add_footer(slide, "Presentation structure: motivation, data, EDA, methods, results, conclusions")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Data Sources And Sample")
    add_callout(
        slide,
        "Dataset 1",
        "AESO Hourly Metered Volumes and Pool Price and AIL:\nHourly pool price, system load, and import-export flows.",
        Inches(0.75),
        Inches(1.55),
        Inches(5.9),
        Inches(1.45),
    )
    add_callout(
        slide,
        "Dataset 2",
        "AESO Historical Generation Data (CSD):\nHourly generation aggregated by fuel type, including gas, wind, solar, hydro, and coal.",
        Inches(0.75),
        Inches(3.25),
        Inches(5.9),
        Inches(1.6),
        fill_color=LIGHT_GRAY,
    )
    add_callout(
        slide,
        "Engineered features",
        "Net load, reserve-style tightness ratios, fuel shares, lagged price terms, lagged spike indicators, and calendar controls.",
        Inches(0.75),
        Inches(5.15),
        Inches(5.9),
        Inches(1.25),
        fill_color=RED_LIGHT,
    )
    add_callout(
        slide,
        "Merged sample",
        f"{len(df):,} hourly observations\n{df['datetime'].min():%Y-%m-%d %H:%M} to {df['datetime'].max():%Y-%m-%d %H:%M}\nMountain Prevailing Time alignment",
        Inches(7.05),
        Inches(1.7),
        Inches(5.4),
        Inches(1.65),
    )
    add_callout(
        slide,
        "Class imbalance",
        f"Overall spike_lead_2 rate = {100 * spike_rate:.2f}%\nTrain = 12.37%\nValidation = 6.15%\nTest = 1.79%",
        Inches(7.05),
        Inches(3.65),
        Inches(5.4),
        Inches(1.65),
        fill_color=LIGHT_GRAY,
    )
    add_callout(
        slide,
        "Why time-based splits",
        "Train: 2020-01-02 to 2023-11-05\nValidation: 2023-11-06 to 2024-12-11\nTest: 2024-12-12 to 2025-07-30",
        Inches(7.05),
        Inches(5.55),
        Inches(5.4),
        Inches(1.1),
        fill_color=RED_LIGHT,
    )
    add_footer(slide, "Source: AESO datasets merged at the hourly level")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "EDA: Price Distribution And Timing")
    slide.shapes.add_picture(str(hist_path), Inches(0.75), Inches(1.45), width=Inches(5.8))
    slide.shapes.add_picture(str(heatmap_path), Inches(6.8), Inches(1.45), width=Inches(5.8))
    add_callout(
        slide,
        "Main EDA takeaways",
        f"Pool prices are highly right-skewed: median = CAD {price_stats['50%']:.2f}/MWh, 95th percentile = CAD {price_stats['95%']:.2f}/MWh, and maximum = CAD {price_stats['max']:.2f}/MWh. Average prices peak in hours {int(avg_by_hour.index[0])}, {int(avg_by_hour.index[1])}, and {int(avg_by_hour.index[2])}; the highest-price months are {int(avg_by_month.index[0])}, {int(avg_by_month.index[1])}, and {int(avg_by_month.index[2])}.",
        Inches(0.9),
        Inches(6.05),
        Inches(11.7),
        Inches(0.75),
        fill_color=LIGHT_GRAY,
    )
    add_footer(slide, "Left: histogram with spike threshold. Right: average price heatmap by month and hour.")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "EDA: System Conditions Before Spikes")
    rows = [
        ("Current pool price", "CAD/MWh"),
        ("Current demand", "MW"),
        ("Net load", "MW"),
        ("Wind output", "MW"),
        ("Renewables share", "%"),
        ("Reserve margin", "ratio"),
    ]
    table = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.75), Inches(1.55), Inches(7.2), Inches(4.1)).table
    headers = ["Variable", "Unit", "No future spike", "Future spike at t+2"]
    for j, text in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED
        para = cell.text_frame.paragraphs[0]
        para.font.name = "Arial"
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = WHITE

    name_to_col = {
        "Current pool price": "ACTUAL_POOL_PRICE",
        "Current demand": "ACTUAL_AIL",
        "Net load": "net_load",
        "Wind output": "wind_total",
        "Renewables share": "renewables_share",
        "Reserve margin": "reserve_margin",
    }
    for i, (label, unit) in enumerate(rows, start=1):
        no_spike = pre_spike.loc[0, name_to_col[label]]
        yes_spike = pre_spike.loc[1, name_to_col[label]]
        values = [label, unit, format_value(label, no_spike), format_value(label, yes_spike)]
        for j, text in enumerate(values):
            cell = table.cell(i, j)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_GRAY
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Arial"
            para.font.size = Pt(10.5)
            para.font.color.rgb = CHARCOAL

    add_callout(
        slide,
        "Interpretation",
        "Hours followed by spikes are already tighter at time t: current prices and load are higher, while wind output, renewable share, and reserve margin are lower.",
        Inches(8.3),
        Inches(1.8),
        Inches(4.2),
        Inches(1.45),
    )
    add_callout(
        slide,
        "Additional signal",
        f"net_load_3h_change shifts from {df.groupby('spike_lead_2')['net_load_3h_change'].mean().loc[0]:.2f} MW in non-spike rows to {df.groupby('spike_lead_2')['net_load_3h_change'].mean().loc[1]:.2f} MW before future spikes.",
        Inches(8.3),
        Inches(3.65),
        Inches(4.2),
        Inches(1.15),
        fill_color=LIGHT_GRAY,
    )
    add_callout(
        slide,
        "Economic reading",
        "The EDA is consistent with a scarcity narrative: tighter supply-demand conditions and weaker renewable output tend to precede short-run spikes.",
        Inches(8.3),
        Inches(5.1),
        Inches(4.2),
        Inches(1.3),
        fill_color=RED_LIGHT,
    )
    add_footer(slide, "Means are computed by future spike status using spike_lead_2")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Prediction Task, Baselines, And Neural Models")
    add_callout(
        slide,
        "Evaluation design",
        "Target: spike_lead_2\nPrimary metric: F1-score\nThresholds tuned on validation only\nUntouched test period reserved for final evaluation",
        Inches(0.75),
        Inches(1.55),
        Inches(3.6),
        Inches(1.65),
    )
    add_callout(
        slide,
        "Baselines",
        "Naive persistence predicts that the current spike state persists.\nSARIMAX treats price as I(1), uses d = 1, Fourier seasonality, and walk-forward forecasting to avoid leakage.",
        Inches(0.75),
        Inches(3.55),
        Inches(3.6),
        Inches(1.9),
        fill_color=LIGHT_GRAY,
    )
    add_callout(
        slide,
        "MLP",
        "Simplest neural benchmark.\nFlat feature vector with contemporaneous market conditions plus lagged price and spike variables.",
        Inches(4.65),
        Inches(1.6),
        Inches(2.6),
        Inches(2.1),
    )
    add_callout(
        slide,
        "LSTM",
        "Sequential model designed to learn temporal dependence over recent hourly windows without manually expanding all lags.",
        Inches(7.45),
        Inches(1.6),
        Inches(2.6),
        Inches(2.1),
        fill_color=LIGHT_GRAY,
    )
    add_callout(
        slide,
        "CNN",
        "1D convolutional network that detects local ramp and regime patterns in the recent sequence of market conditions.",
        Inches(10.25),
        Inches(1.6),
        Inches(2.35),
        Inches(2.1),
        fill_color=RED_LIGHT,
    )
    add_callout(
        slide,
        "Why compare all three on one slide",
        "The modeling strategy intentionally moves from a simple MLP benchmark to richer temporal architectures. This makes it easier to judge whether added sequential structure actually improves spike detection.",
        Inches(4.65),
        Inches(4.15),
        Inches(7.95),
        Inches(1.45),
        fill_color=LIGHT_GRAY,
    )
    add_footer(slide, "Same split logic and validation-threshold tuning were used across models")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Results Comparison")
    slide.shapes.add_picture(str(model_chart_path), Inches(0.75), Inches(1.55), width=Inches(5.8))
    result_table = slide.shapes.add_table(6, 4, Inches(6.9), Inches(1.7), Inches(5.3), Inches(3.8)).table
    result_headers = ["Model", "F1", "Precision", "Recall"]
    for j, text in enumerate(result_headers):
        cell = result_table.cell(0, j)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED
        para = cell.text_frame.paragraphs[0]
        para.font.name = "Arial"
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = WHITE

    metrics_rows = [
        ("Naive", baseline_rows.loc["Naive", "test_f1"], baseline_rows.loc["Naive", "test_precision"], baseline_rows.loc["Naive", "test_recall"]),
        ("SARIMAX", baseline_rows.loc["SARIMAX + Fourier", "test_f1"], baseline_rows.loc["SARIMAX + Fourier", "test_precision"], baseline_rows.loc["SARIMAX + Fourier", "test_recall"]),
        ("MLP", mlp["test"]["f1"], mlp["test"]["precision"], mlp["test"]["recall"]),
        ("LSTM", lstm["test"]["f1"], lstm["test"]["precision"], lstm["test"]["recall"]),
        ("CNN", cnn["f1"], cnn["precision"], cnn["recall"]),
    ]
    for i, row in enumerate(metrics_rows, start=1):
        for j, value in enumerate(row):
            cell = result_table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_GRAY
            cell.text = value if j == 0 else f"{float(value):.3f}"
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Arial"
            para.font.size = Pt(10.5)
            para.font.color.rgb = CHARCOAL
            if i == 5:
                para.font.bold = True

    add_callout(
        slide,
        "Headline result",
        "The CNN achieved the strongest held-out performance (test F1 = 0.416), while the naive baseline remained surprisingly competitive at 0.384.",
        Inches(6.9),
        Inches(5.85),
        Inches(5.3),
        Inches(0.8),
        fill_color=RED_LIGHT,
    )
    add_footer(slide, "All metrics are from the untouched test period")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Interpretation")
    add_callout(
        slide,
        "Finding 1",
        "The CNN is the only neural model that clearly improved on the naive persistence baseline in held-out F1.",
        Inches(0.75),
        Inches(1.7),
        Inches(3.95),
        Inches(1.55),
    )
    add_callout(
        slide,
        "Finding 2",
        "The LSTM and naive baseline were very close. This implies short-run persistence already captures a large share of the available predictive signal.",
        Inches(4.95),
        Inches(1.7),
        Inches(3.95),
        Inches(1.55),
        fill_color=LIGHT_GRAY,
    )
    add_callout(
        slide,
        "Finding 3",
        "The MLP remained a useful benchmark because it showed how much can be learned from lagged variables alone before moving to sequence models.",
        Inches(9.15),
        Inches(1.7),
        Inches(3.15),
        Inches(1.55),
        fill_color=RED_LIGHT,
    )
    add_bullets(
        slide,
        [
            "Across the neural models, the same drivers repeatedly surfaced: ACTUAL_POOL_PRICE, wind_total, solar_total, and net_load_3h_change.",
            "Those variables align with the EDA and support an economic interpretation based on system tightness and renewable variability.",
            "The baseline comparison is substantively important: without it, the LSTM would appear stronger than it really is.",
        ],
        Inches(0.9),
        Inches(3.9),
        Inches(11.5),
        Inches(2.1),
        font_size=18,
    )
    add_footer(slide, "Interpretation combines model comparison, feature importance, and system-level economics")

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Conclusions And Future Work")
    add_bullets(
        slide,
        [
            "Public AESO data contains meaningful signal for short-horizon Alberta electricity price spike prediction.",
            "The CNN produced the strongest overall result, but strong simple baselines show the task is harder than raw neural scores alone suggest.",
            "The main practical lesson is that richer temporal structure helps, but model gains remain modest in a changing and imbalanced market.",
        ],
        Inches(0.8),
        Inches(1.55),
        Inches(7.2),
        Inches(3.1),
    )
    add_callout(
        slide,
        "Future work",
        "Add outage information, weather forecasts, probability calibration, and regime-specific models to improve robustness under changing system conditions.",
        Inches(8.35),
        Inches(1.85),
        Inches(4.1),
        Inches(1.6),
    )
    add_callout(
        slide,
        "Presentation takeaway",
        "For this project, the CNN delivered the best balance of precision and recall, while the MLP provided a simple benchmark and the LSTM showed that sequential memory alone was not enough to dominate persistence.",
        Inches(8.35),
        Inches(4.05),
        Inches(4.1),
        Inches(1.9),
        fill_color=LIGHT_GRAY,
    )
    add_footer(slide, "End of presentation")

    prs.save(OUTPUT_PPTX)
    print(f"Saved deck to {OUTPUT_PPTX}")


if __name__ == "__main__":
    build_deck()
